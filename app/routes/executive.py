"""
Executive Dashboard routes
"""

from datetime import datetime, timedelta

from flask import Blueprint, redirect, render_template, request, session, url_for
from flask_login import login_required
from sqlalchemy import func

from app.extensions import db
from app.models import KPI, Contribution, GovernanceBody, InitiativeSystemLink, KPIValueTypeConfig, System

bp = Blueprint("executive", __name__, url_prefix="/executive")


@bp.route("/dashboard")
@login_required
def dashboard():
    """
    Executive Dashboard - High-level KPI health and strategic overview.

    Shows traffic light status, governance body performance, top/bottom performers,
    and key insights for executive decision-making.
    """
    org_id = session.get("organization_id")
    if not org_id:
        return redirect(url_for("auth.login"))

    # Date range selector (default: 30 days)
    days = request.args.get("days", 30, type=int)
    today = datetime.now().date()
    period_start = today - timedelta(days=days)

    # === GET ALL KPIs WITH STATUS ===

    kpis = (
        KPI.query.join(InitiativeSystemLink)
        .join(System)
        .filter(System.organization_id == org_id)
        .order_by(KPI.name)
        .all()
    )

    # Calculate status for each KPI
    kpi_statuses = []
    for kpi in kpis:
        status_data = kpi.get_status()
        kpi_statuses.append(
            {
                "kpi": kpi,
                "status": status_data["status"],
                "reason": status_data["reason"],
                "details": status_data["details"],
            }
        )

    # === KEY METRICS ===

    total_kpis = len(kpi_statuses)
    active_kpis = sum(1 for k in kpi_statuses if not k["kpi"].is_archived)
    green_count = sum(1 for k in kpi_statuses if k["status"] == "green")
    yellow_count = sum(1 for k in kpi_statuses if k["status"] == "yellow")
    red_count = sum(1 for k in kpi_statuses if k["status"] == "red")

    # Calculate percentages
    green_pct = (green_count / active_kpis * 100) if active_kpis > 0 else 0
    yellow_pct = (yellow_count / active_kpis * 100) if active_kpis > 0 else 0
    red_pct = (red_count / active_kpis * 100) if active_kpis > 0 else 0

    # Calculate "On Target" percentage (green + yellow achieving > 60%)
    on_target_count = green_count + sum(
        1 for k in kpi_statuses if k["status"] == "yellow" and k["details"].get("target_achievement_pct", 0) >= 60
    )
    on_target_pct = (on_target_count / active_kpis * 100) if active_kpis > 0 else 0

    # === TREND CALCULATION ===

    # Count contributions in current period
    current_contributions = (
        db.session.query(func.count(Contribution.id))
        .join(KPIValueTypeConfig, Contribution.kpi_value_type_config_id == KPIValueTypeConfig.id)
        .join(KPI, KPIValueTypeConfig.kpi_id == KPI.id)
        .join(InitiativeSystemLink, KPI.initiative_system_link_id == InitiativeSystemLink.id)
        .join(System, InitiativeSystemLink.system_id == System.id)
        .filter(System.organization_id == org_id, Contribution.created_at >= period_start)
        .scalar()
    )

    # Count contributions in previous period (for trend comparison)
    previous_period_start = period_start - timedelta(days=days)
    previous_contributions = (
        db.session.query(func.count(Contribution.id))
        .join(KPIValueTypeConfig, Contribution.kpi_value_type_config_id == KPIValueTypeConfig.id)
        .join(KPI, KPIValueTypeConfig.kpi_id == KPI.id)
        .join(InitiativeSystemLink, KPI.initiative_system_link_id == InitiativeSystemLink.id)
        .join(System, InitiativeSystemLink.system_id == System.id)
        .filter(
            System.organization_id == org_id,
            Contribution.created_at >= previous_period_start,
            Contribution.created_at < period_start,
        )
        .scalar()
    )

    # Calculate trend
    if previous_contributions > 0:
        activity_trend_pct = ((current_contributions - previous_contributions) / previous_contributions) * 100
    else:
        activity_trend_pct = 100 if current_contributions > 0 else 0

    # === GOVERNANCE BODY PERFORMANCE ===

    governance_bodies = (
        GovernanceBody.query.filter_by(organization_id=org_id).order_by(GovernanceBody.display_order).all()
    )

    gov_body_performance = []
    for gb in governance_bodies:
        # Get KPIs for this governance body
        gb_kpi_ids = [link.kpi_id for link in gb.kpi_links]
        gb_kpis = [k for k in kpi_statuses if k["kpi"].id in gb_kpi_ids]

        if gb_kpis:
            gb_green = sum(1 for k in gb_kpis if k["status"] == "green")
            gb_total = len(gb_kpis)
            gb_pct = (gb_green / gb_total * 100) if gb_total > 0 else 0

            # Overall status
            gb_red = sum(1 for k in gb_kpis if k["status"] == "red")
            if gb_pct >= 75:
                gb_status = "green"
            elif gb_red > 0 or gb_pct < 50:
                gb_status = "red"
            else:
                gb_status = "yellow"

            gov_body_performance.append(
                {
                    "name": gb.name,
                    "color": gb.color,
                    "performance_pct": round(gb_pct, 0),
                    "status": gb_status,
                    "kpi_count": gb_total,
                }
            )

    # === TOP PERFORMERS ===

    # KPIs with green status and target achievement > 100%
    top_performers = []
    for k in kpi_statuses:
        if k["status"] == "green" and k["details"].get("target_achievement_pct"):
            achievement_pct = k["details"]["target_achievement_pct"]
            if achievement_pct >= 100:
                top_performers.append(
                    {
                        "kpi": k["kpi"],
                        "achievement_pct": achievement_pct,
                        "system": (
                            k["kpi"].initiative_system_link.system.name if k["kpi"].initiative_system_link else "N/A"
                        ),
                    }
                )

    top_performers.sort(key=lambda x: x["achievement_pct"], reverse=True)
    top_performers = top_performers[:5]  # Top 5

    # === NEEDS ATTENTION ===

    # Red KPIs or yellow KPIs with low achievement
    needs_attention = []
    for k in kpi_statuses:
        if k["status"] == "red" or (k["status"] == "yellow" and k["details"].get("target_achievement_pct", 100) < 70):
            needs_attention.append(
                {
                    "kpi": k["kpi"],
                    "status": k["status"],
                    "reason": k["reason"],
                    "achievement_pct": k["details"].get("target_achievement_pct"),
                    "days_since_activity": k["details"].get("days_since_activity"),
                    "system": k["kpi"].initiative_system_link.system.name if k["kpi"].initiative_system_link else "N/A",
                }
            )

    needs_attention.sort(
        key=lambda x: (0 if x["status"] == "red" else 1, x.get("days_since_activity") or 999), reverse=False
    )
    needs_attention = needs_attention[:10]  # Top 10

    # === ACTIVITY TREND (Weekly) ===

    # Get contributions by week for chart
    weeks = []
    for i in range(4):
        week_end = today - timedelta(days=i * 7)
        week_start = week_end - timedelta(days=7)

        week_contributions = (
            db.session.query(func.count(Contribution.id))
            .join(KPIValueTypeConfig, Contribution.kpi_value_type_config_id == KPIValueTypeConfig.id)
            .join(KPI, KPIValueTypeConfig.kpi_id == KPI.id)
            .join(InitiativeSystemLink, KPI.initiative_system_link_id == InitiativeSystemLink.id)
            .join(System, InitiativeSystemLink.system_id == System.id)
            .filter(
                System.organization_id == org_id,
                Contribution.created_at >= week_start,
                Contribution.created_at < week_end,
            )
            .scalar()
        )

        weeks.append({"label": f"Week {4-i}", "count": week_contributions})

    weeks.reverse()  # Chronological order

    # === KEY INSIGHTS ===

    insights = []

    # Positive insight
    if green_pct >= 70:
        insights.append(
            {
                "type": "success",
                "title": "Strong Performance",
                "message": f"{green_count} KPIs ({green_pct:.0f}%) are on track with recent activity and good consensus.",
            }
        )

    # Warning insight
    if red_count > 0:
        stale_count = sum(1 for k in needs_attention if k.get("days_since_activity") and k["days_since_activity"] >= 30)
        if stale_count > 0:
            insights.append(
                {
                    "type": "warning",
                    "title": "Action Required",
                    "message": f"{stale_count} KPIs have had no activity in 30+ days. Immediate attention needed.",
                }
            )

    # Trend insight
    if activity_trend_pct > 0:
        insights.append(
            {
                "type": "info",
                "title": "Positive Momentum",
                "message": f"Activity increased {activity_trend_pct:.0f}% compared to previous period. Keep up the momentum!",
            }
        )
    elif activity_trend_pct < -20:
        insights.append(
            {
                "type": "warning",
                "title": "Declining Activity",
                "message": f"Activity decreased {abs(activity_trend_pct):.0f}% compared to previous period. Consider re-engagement.",
            }
        )

    # === ALERTS ===

    alerts = []

    # Overdue targets
    overdue_count = sum(1 for k in kpi_statuses if k["details"].get("target_status") == "overdue")
    if overdue_count > 0:
        alerts.append(
            {
                "message": f"{overdue_count} KPI{'s' if overdue_count != 1 else ''} missed their target deadline.",
                "url": "#needs-attention",
            }
        )

    # Approaching deadlines
    approaching_count = sum(1 for k in kpi_statuses if k["details"].get("target_status") == "approaching")
    if approaching_count > 0:
        alerts.append(
            {
                "message": f"{approaching_count} KPI{'s' if approaching_count != 1 else ''} have deadlines within 7 days.",
                "url": "#needs-attention",
            }
        )

    # Stale KPIs
    stale_count = sum(
        1
        for k in kpi_statuses
        if k["details"].get("days_since_activity") is None or k["details"].get("days_since_activity", 0) >= 30
    )
    if stale_count > 0:
        alerts.append(
            {
                "message": f"{stale_count} KPI{'s' if stale_count != 1 else ''} have no activity in the last 30 days.",
                "url": "#needs-attention",
            }
        )

    return render_template(
        "executive/dashboard.html",
        # Metrics
        total_kpis=total_kpis,
        active_kpis=active_kpis,
        on_target_pct=round(on_target_pct, 0),
        needs_attention_count=yellow_count,
        at_risk_count=red_count,
        # Traffic lights
        green_count=green_count,
        yellow_count=yellow_count,
        red_count=red_count,
        green_pct=round(green_pct, 0),
        yellow_pct=round(yellow_pct, 0),
        red_pct=round(red_pct, 0),
        # Trends
        activity_trend_pct=round(activity_trend_pct, 1),
        # Governance
        gov_body_performance=gov_body_performance,
        # Lists
        top_performers=top_performers,
        needs_attention=needs_attention,
        # Charts
        weeks=weeks,
        # Insights & Alerts
        insights=insights,
        alerts=alerts,
        # Period
        days=days,
    )


@bp.route("/export/excel")
@login_required
def export_excel():
    """
    Export executive dashboard to Excel format.
    """
    from io import BytesIO

    from flask import send_file
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    org_id = session.get("organization_id")
    org_name = session.get("organization_name")
    days = request.args.get("days", 30, type=int)

    # Get data (reuse logic from dashboard)
    kpis = KPI.query.join(InitiativeSystemLink).join(System).filter(System.organization_id == org_id).all()

    kpi_statuses = []
    for kpi in kpis:
        status_data = kpi.get_status()
        kpi_statuses.append(
            {
                "kpi": kpi,
                "status": status_data["status"],
                "reason": status_data["reason"],
                "details": status_data["details"],
            }
        )

    # Create workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Executive Summary"

    # Header
    ws["A1"] = "EXECUTIVE DASHBOARD"
    ws["A1"].font = Font(size=18, bold=True)
    ws["A2"] = f"{org_name} - Last {days} Days"
    ws["A2"].font = Font(size=12)
    ws["A3"] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # Summary Stats
    row = 5
    ws[f"A{row}"] = "SUMMARY"
    ws[f"A{row}"].font = Font(size=14, bold=True)

    green_count = sum(1 for k in kpi_statuses if k["status"] == "green")
    yellow_count = sum(1 for k in kpi_statuses if k["status"] == "yellow")
    red_count = sum(1 for k in kpi_statuses if k["status"] == "red")

    row += 2
    ws[f"A{row}"] = "Status"
    ws[f"B{row}"] = "Count"
    ws[f"C{row}"] = "Percentage"
    for cell in [f"A{row}", f"B{row}", f"C{row}"]:
        ws[cell].font = Font(bold=True)
        ws[cell].fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")

    total = len(kpi_statuses)

    row += 1
    ws[f"A{row}"] = "🟢 On Track"
    ws[f"B{row}"] = green_count
    ws[f"C{row}"] = f"{(green_count/total*100):.0f}%" if total > 0 else "0%"
    ws[f"A{row}"].fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")

    row += 1
    ws[f"A{row}"] = "🟡 Needs Attention"
    ws[f"B{row}"] = yellow_count
    ws[f"C{row}"] = f"{(yellow_count/total*100):.0f}%" if total > 0 else "0%"
    ws[f"A{row}"].fill = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")

    row += 1
    ws[f"A{row}"] = "🔴 At Risk"
    ws[f"B{row}"] = red_count
    ws[f"C{row}"] = f"{(red_count/total*100):.0f}%" if total > 0 else "0%"
    ws[f"A{row}"].fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")

    # KPI Details
    row += 3
    ws[f"A{row}"] = "KPI DETAILS"
    ws[f"A{row}"].font = Font(size=14, bold=True)

    row += 2
    headers = ["KPI Name", "System", "Status", "Reason", "Days Since Activity", "Target Achievement"]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=row, column=col, value=header)
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")

    for k in kpi_statuses:
        row += 1
        ws.cell(row=row, column=1, value=k["kpi"].name)
        ws.cell(
            row=row,
            column=2,
            value=k["kpi"].initiative_system_link.system.name if k["kpi"].initiative_system_link else "N/A",
        )
        ws.cell(row=row, column=3, value=k["status"].upper())
        ws.cell(row=row, column=4, value=k["reason"])
        ws.cell(row=row, column=5, value=k["details"].get("days_since_activity", "N/A"))

        target_pct = k["details"].get("target_achievement_pct")
        ws.cell(row=row, column=6, value=f"{target_pct:.1f}%" if target_pct else "N/A")

        # Color code status
        if k["status"] == "green":
            ws.cell(row=row, column=3).fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
        elif k["status"] == "yellow":
            ws.cell(row=row, column=3).fill = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
        else:
            ws.cell(row=row, column=3).fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")

    # Adjust column widths
    ws.column_dimensions["A"].width = 40
    ws.column_dimensions["B"].width = 25
    ws.column_dimensions["C"].width = 15
    ws.column_dimensions["D"].width = 50
    ws.column_dimensions["E"].width = 20
    ws.column_dimensions["F"].width = 20

    # Save to BytesIO
    bio = BytesIO()
    wb.save(bio)
    bio.seek(0)

    filename = f"Executive_Dashboard_{org_name}_{datetime.now().strftime('%Y%m%d')}.xlsx"

    return send_file(
        bio,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        as_attachment=True,
        download_name=filename,
    )


@bp.route("/export/pdf")
@login_required
def export_pdf():
    """
    Export executive dashboard to PDF format.

    Simple HTML-to-PDF using browser print CSS.
    """
    org_id = session.get("organization_id")
    if not org_id:
        return redirect(url_for("auth.login"))

    days = request.args.get("days", 30, type=int)

    # Render the same dashboard template but with print-friendly CSS
    # For now, just render the HTML page - user can use browser "Print to PDF"
    # In future, we can use WeasyPrint or similar library for server-side PDF generation

    return render_template(
        "executive/dashboard_print.html",
        message="Please use your browser's Print function (Ctrl+P / Cmd+P) and select 'Save as PDF'",
        dashboard_url=url_for("executive.dashboard", days=days),
    )


@bp.route("/export/powerpoint")
@login_required
def export_powerpoint():
    """
    Export executive dashboard to PowerPoint format.
    """
    from io import BytesIO

    from flask import send_file
    from pptx import Presentation
    from pptx.util import Inches, Pt

    org_id = session.get("organization_id")
    org_name = session.get("organization_name")

    # Get data
    kpis = KPI.query.join(InitiativeSystemLink).join(System).filter(System.organization_id == org_id).all()

    kpi_statuses = []
    for kpi in kpis:
        status_data = kpi.get_status()
        kpi_statuses.append(
            {
                "kpi": kpi,
                "status": status_data["status"],
                "reason": status_data["reason"],
                "details": status_data["details"],
            }
        )

    green_count = sum(1 for k in kpi_statuses if k["status"] == "green")
    yellow_count = sum(1 for k in kpi_statuses if k["status"] == "yellow")
    red_count = sum(1 for k in kpi_statuses if k["status"] == "red")
    total = len(kpi_statuses)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Executive Dashboard"
    subtitle.text = f"{org_name}\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"

    # Slide 2: KPI Health Status
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "KPI Health Status"

    # Add text boxes for traffic lights
    left = Inches(1)
    top = Inches(2)
    width = Inches(2.5)
    height = Inches(1.5)

    # Green
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"🟢 On Track\n{green_count} KPIs\n{(green_count/total*100):.0f}%" if total > 0 else "0%"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    # Yellow
    txBox = slide.shapes.add_textbox(left + Inches(3), top, width, height)
    tf = txBox.text_frame
    tf.text = f"🟡 Needs Attention\n{yellow_count} KPIs\n{(yellow_count/total*100):.0f}%" if total > 0 else "0%"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    # Red
    txBox = slide.shapes.add_textbox(left + Inches(6), top, width, height)
    tf = txBox.text_frame
    tf.text = f"🔴 At Risk\n{red_count} KPIs\n{(red_count/total*100):.0f}%" if total > 0 else "0%"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    # Slide 3: Needs Attention
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "KPIs Requiring Attention"

    needs_attention = [k for k in kpi_statuses if k["status"] in ["red", "yellow"]][:10]

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    if needs_attention:
        for k in needs_attention:
            p = tf.add_paragraph()
            status_icon = "🔴" if k["status"] == "red" else "🟡"
            p.text = f"{status_icon} {k['kpi'].name}: {k['reason']}"
            p.font.size = Pt(14)
            p.space_after = Pt(10)
    else:
        p = tf.paragraphs[0]
        p.text = "✅ All KPIs are on track!"
        p.font.size = Pt(18)

    # Save to BytesIO
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    filename = f"Executive_Dashboard_{org_name}_{datetime.now().strftime('%Y%m%d')}.pptx"

    return send_file(
        bio,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )
